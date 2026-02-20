"""
build_presentation.py — Generate Eclipse Scalper technical PowerPoint deck.

Usage:
    python -m tools.build_presentation
    python -m tools.build_presentation --out reports/eclipse_scalper_deck.pptx
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Colour palette — dark engineering theme
# ---------------------------------------------------------------------------
C_BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_BG_MID       = RGBColor(0x1A, 0x2B, 0x3C)   # section bg
C_BG_SLIDE     = RGBColor(0x0F, 0x1F, 0x30)   # slide bg
C_ACCENT       = RGBColor(0x00, 0xD4, 0xFF)   # cyan accent
C_ACCENT2      = RGBColor(0x00, 0xFF, 0xB3)   # mint accent
C_WARNING      = RGBColor(0xFF, 0xA5, 0x00)   # amber
C_DANGER       = RGBColor(0xFF, 0x3C, 0x3C)   # red
C_SUCCESS      = RGBColor(0x00, 0xC8, 0x5A)   # green
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GREY   = RGBColor(0xB0, 0xC4, 0xD8)
C_MID_GREY     = RGBColor(0x60, 0x80, 0x9A)
C_TITLE_BAR    = RGBColor(0x00, 0x8B, 0xCC)   # title bar blue
C_SECTION      = RGBColor(0x1A, 0x3A, 0x5C)   # section divider

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _solid_fill(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _add_rect(slide, left, top, width, height, rgb: RGBColor, transparency: int = 0):
    shape = slide.shapes.add_shape(1, left, top, width, height)   # MSO_SHAPE_TYPE.RECTANGLE=1
    _solid_fill(shape, rgb)
    shape.line.fill.background()
    if transparency:
        shape.fill.fore_color.theme_color = None  # ensure rgb mode
    return shape


def _set_slide_bg(slide, rgb: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _tf(shape):
    return shape.text_frame


def _para(tf, text: str, size: int, bold: bool, colour: RGBColor,
          align=PP_ALIGN.LEFT, space_before: int = 0, space_after: int = 0,
          italic: bool = False, level: int = 0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    return p


def _add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb


def _clear_tf(tf):
    """Remove the default empty paragraph that pptx adds."""
    # keep at least one; we will overwrite
    pass


def _notes(slide, text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)


def _visual_box(slide, left, top, width, height, label: str):
    """Grey placeholder box describing the intended visual."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x15, 0x30, 0x45)
    box.line.color.rgb = C_ACCENT
    box.line.width = Pt(1)
    tf2 = box.text_frame
    tf2.word_wrap = True
    tf2.auto_size = None
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[ VISUAL ]"
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = C_ACCENT
    p2 = tf2.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(7)
    run2.font.italic = True
    run2.font.color.rgb = C_LIGHT_GREY


# ---------------------------------------------------------------------------
# Slide builder functions
# ---------------------------------------------------------------------------

def _title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    _set_slide_bg(slide, C_BG_DARK)

    # Top accent bar
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT)

    # Left side accent stripe
    _add_rect(slide, 0, Inches(0.08), Inches(0.06), SLIDE_H - Inches(0.08), C_TITLE_BAR)

    # Main title block
    tb = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(12.5), Inches(2.2))
    tf = _tf(tb)
    _para(tf, "ECLIPSE SCALPER", 42, True, C_ACCENT, PP_ALIGN.LEFT, 0, 4)
    _para(tf, "Microstructure Alpha Discovery, Validation,", 26, False, C_WHITE, PP_ALIGN.LEFT, 0, 0)
    _para(tf, "and Execution Realism", 26, False, C_WHITE, PP_ALIGN.LEFT, 0, 8)

    # Separator line (thin rect)
    _add_rect(slide, Inches(0.5), Inches(3.55), Inches(8), Inches(0.03), C_ACCENT2)

    # Subtitle
    tb2 = _add_textbox(slide, Inches(0.5), Inches(3.7), Inches(12), Inches(0.6))
    tf2 = _tf(tb2)
    _para(tf2, "A Systems Engineering & Quantitative Research Deep-Dive", 16, False, C_LIGHT_GREY, PP_ALIGN.LEFT)

    # Meta info
    tb3 = _add_textbox(slide, Inches(0.5), Inches(4.4), Inches(9), Inches(1.5))
    tf3 = _tf(tb3)
    _para(tf3, "Internal Engineering Review  |  Quantitative Research Division", 12, False, C_MID_GREY, PP_ALIGN.LEFT, 0, 4)
    _para(tf3, "Branch: feat/reliability-gate-automation  |  Date: 2026-02-20", 11, False, C_MID_GREY, PP_ALIGN.LEFT, 0, 4)
    _para(tf3, "79 tests passing  |  Python 3.13  |  python-pptx generated", 10, False, C_MID_GREY, PP_ALIGN.LEFT)

    # Bottom bar
    _add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), C_ACCENT)

    # Pipeline label
    tb4 = _add_textbox(slide, Inches(9.5), Inches(4.0), Inches(3.5), Inches(2.8))
    tf4 = _tf(tb4)
    for line in [
        "signal",
        "  ↓",
        "entry_loop",
        "  ↓",
        "order_router",
        "  ↓",
        "exchange",
        "  ↓",
        "reconcile",
    ]:
        _para(tf4, line, 9, False, C_ACCENT if "↓" not in line else C_MID_GREY, PP_ALIGN.LEFT, 1, 1)

    _notes(slide,
        "This presentation documents the complete research and engineering stack of Eclipse Scalper — "
        "an async Python cryptocurrency futures trading system targeting Binance perpetuals via ccxt. "
        "The presentation is structured for audiences with backgrounds in quantitative finance, "
        "distributed systems, and algorithmic execution. It is intentionally sober: where the system "
        "finds negative results, those are presented as such. The absence of deployed alpha is treated "
        "as a feature of the measurement system, not a defect."
    )
    return slide


def _section_divider(prs: Presentation, number: str, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_SECTION)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT)
    _add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), C_ACCENT)

    # Number
    tb0 = _add_textbox(slide, Inches(0.6), Inches(1.8), Inches(1.5), Inches(1.5))
    _para(_tf(tb0), number, 72, True, C_ACCENT, PP_ALIGN.LEFT)

    # Titles
    tb = _add_textbox(slide, Inches(2.2), Inches(2.2), Inches(10.5), Inches(2.5))
    tf = _tf(tb)
    _para(tf, title, 34, True, C_WHITE, PP_ALIGN.LEFT, 0, 6)
    _para(tf, subtitle, 16, False, C_LIGHT_GREY, PP_ALIGN.LEFT)

    _add_rect(slide, Inches(2.2), Inches(2.1), Inches(0.04), Inches(2.2), C_ACCENT2)
    _notes(slide, f"Section divider: {title} — {subtitle}")
    return slide


def _content_slide(
    prs: Presentation,
    slide_id: str,
    title: str,
    bullets: list[tuple[str, int]],   # (text, level)  level 0=main 1=sub
    notes_text: str,
    visual_label: str = "",
    has_visual: bool = True,
    visual_right: bool = True,
):
    """Generic content slide. If has_visual, right 40% is a visual placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_SLIDE)

    # Top bar
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT)

    # Title bar
    title_bar = _add_rect(slide, 0, Inches(0.06), SLIDE_W, Inches(0.72), C_BG_MID)
    tb_title = _add_textbox(slide, Inches(0.3), Inches(0.08), Inches(11.8), Inches(0.65))
    tf_t = _tf(tb_title)
    _para(tf_t, title, 20, True, C_WHITE, PP_ALIGN.LEFT)

    # Slide ID chip
    tb_id = _add_textbox(slide, Inches(12.2), Inches(0.1), Inches(1.0), Inches(0.5))
    _para(_tf(tb_id), slide_id, 8, False, C_ACCENT, PP_ALIGN.RIGHT)

    content_left  = Inches(0.3)
    content_top   = Inches(0.88)
    content_width = Inches(7.6) if (has_visual and visual_right) else Inches(12.7)
    content_height = Inches(6.3)

    tb_body = _add_textbox(slide, content_left, content_top, content_width, content_height)
    tf_b = _tf(tb_body)
    tf_b.word_wrap = True

    first = True
    for text, level in bullets:
        if first:
            p = tf_b.paragraphs[0]
            first = False
        else:
            p = tf_b.add_paragraph()
        p.level = level
        p.space_before = Pt(3 if level == 0 else 1)
        p.space_after = Pt(2)
        p.alignment = PP_ALIGN.LEFT

        bullet_char = "▸ " if level == 0 else "  · "
        bullet_colour = C_ACCENT if level == 0 else C_ACCENT2
        text_colour = C_WHITE if level == 0 else C_LIGHT_GREY
        font_size = 13 if level == 0 else 11

        run_b = p.add_run()
        run_b.text = bullet_char
        run_b.font.size = Pt(font_size)
        run_b.font.bold = (level == 0)
        run_b.font.color.rgb = bullet_colour

        run_t = p.add_run()
        run_t.text = text
        run_t.font.size = Pt(font_size)
        run_t.font.bold = False
        run_t.font.color.rgb = text_colour

    if has_visual and visual_right:
        _visual_box(
            slide,
            Inches(8.1), Inches(0.88),
            Inches(5.0), Inches(6.3),
            visual_label,
        )

    # Bottom bar
    _add_rect(slide, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), C_TITLE_BAR)

    _notes(slide, notes_text)
    return slide


def _code_slide(prs: Presentation, slide_id: str, title: str, code_lines: list[str], notes_text: str):
    """Slide for code/formula content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_SLIDE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT)
    _add_rect(slide, 0, Inches(0.06), SLIDE_W, Inches(0.72), C_BG_MID)

    tb_title = _add_textbox(slide, Inches(0.3), Inches(0.08), Inches(11.8), Inches(0.65))
    _para(_tf(tb_title), title, 20, True, C_WHITE, PP_ALIGN.LEFT)

    tb_id = _add_textbox(slide, Inches(12.2), Inches(0.1), Inches(1.0), Inches(0.5))
    _para(_tf(tb_id), slide_id, 8, False, C_ACCENT, PP_ALIGN.RIGHT)

    # Code box
    code_box = _add_rect(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(6.15), RGBColor(0x05, 0x12, 0x1E))
    code_box.line.color.rgb = C_ACCENT
    code_box.line.width = Pt(0.75)

    tb_code = _add_textbox(slide, Inches(0.65), Inches(1.05), Inches(12.0), Inches(5.9))
    tf_c = _tf(tb_code)
    tf_c.word_wrap = False

    first = True
    for line in code_lines:
        if first:
            p = tf_c.paragraphs[0]
            first = False
        else:
            p = tf_c.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line

        # colour keywords
        colour = C_LIGHT_GREY
        if line.strip().startswith("#"):
            colour = C_MID_GREY
        elif any(kw in line for kw in ["def ", "class ", "import ", "from "]):
            colour = C_ACCENT2
        elif any(kw in line for kw in ["=", "return", "assert", "if ", "else"]):
            colour = C_WARNING
        run.font.size = Pt(9.5)
        run.font.color.rgb = colour
        run.font.name = "Courier New"

    _add_rect(slide, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), C_TITLE_BAR)
    _notes(slide, notes_text)
    return slide


def _two_col_slide(prs, slide_id, title, left_items, right_items, notes_text,
                   left_header="", right_header=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_SLIDE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT)
    _add_rect(slide, 0, Inches(0.06), SLIDE_W, Inches(0.72), C_BG_MID)

    tb_title = _add_textbox(slide, Inches(0.3), Inches(0.08), Inches(11.8), Inches(0.65))
    _para(_tf(tb_title), title, 20, True, C_WHITE, PP_ALIGN.LEFT)

    tb_id = _add_textbox(slide, Inches(12.2), Inches(0.1), Inches(1.0), Inches(0.5))
    _para(_tf(tb_id), slide_id, 8, False, C_ACCENT, PP_ALIGN.RIGHT)

    # divider
    _add_rect(slide, Inches(6.55), Inches(0.9), Inches(0.03), Inches(6.3), C_MID_GREY)

    for col_idx, (header, items) in enumerate([(left_header, left_items), (right_header, right_items)]):
        x = Inches(0.35) if col_idx == 0 else Inches(6.7)
        w = Inches(6.1)
        tb_h = _add_textbox(slide, x, Inches(0.9), w, Inches(0.4))
        _para(_tf(tb_h), header, 11, True, C_ACCENT2, PP_ALIGN.LEFT)

        tb_body = _add_textbox(slide, x, Inches(1.35), w, Inches(5.8))
        tf_b = _tf(tb_body)
        tf_b.word_wrap = True
        first = True
        for text in items:
            if first:
                p = tf_b.paragraphs[0]
                first = False
            else:
                p = tf_b.add_paragraph()
            p.space_before = Pt(3)
            run_b = p.add_run()
            run_b.text = "▸ "
            run_b.font.size = Pt(11)
            run_b.font.bold = True
            run_b.font.color.rgb = C_ACCENT
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(11)
            run_t.font.color.rgb = C_WHITE

    _add_rect(slide, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), C_TITLE_BAR)
    _notes(slide, notes_text)
    return slide


def _formula_slide(prs, slide_id, title, formulas, explanation_bullets, notes_text):
    """Slide with formula box on left, explanation on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_SLIDE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT)
    _add_rect(slide, 0, Inches(0.06), SLIDE_W, Inches(0.72), C_BG_MID)

    tb_title = _add_textbox(slide, Inches(0.3), Inches(0.08), Inches(11.8), Inches(0.65))
    _para(_tf(tb_title), title, 20, True, C_WHITE, PP_ALIGN.LEFT)

    tb_id = _add_textbox(slide, Inches(12.2), Inches(0.1), Inches(1.0), Inches(0.5))
    _para(_tf(tb_id), slide_id, 8, False, C_ACCENT, PP_ALIGN.RIGHT)

    # Formula box
    fbox = _add_rect(slide, Inches(0.4), Inches(0.95), Inches(6.0), Inches(6.15), RGBColor(0x05, 0x12, 0x1E))
    fbox.line.color.rgb = C_ACCENT2
    fbox.line.width = Pt(0.75)

    tb_f = _add_textbox(slide, Inches(0.6), Inches(1.05), Inches(5.6), Inches(5.9))
    tf_f = _tf(tb_f)
    tf_f.word_wrap = True
    first = True
    for f in formulas:
        if first:
            p = tf_f.paragraphs[0]
            first = False
        else:
            p = tf_f.add_paragraph()
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f
        run.font.size = Pt(11)
        run.font.name = "Courier New"
        run.font.color.rgb = C_ACCENT2

    # Explanation
    tb_e = _add_textbox(slide, Inches(6.7), Inches(0.95), Inches(6.4), Inches(6.15))
    tf_e = _tf(tb_e)
    tf_e.word_wrap = True
    first = True
    for text, level in explanation_bullets:
        if first:
            p = tf_e.paragraphs[0]
            first = False
        else:
            p = tf_e.add_paragraph()
        p.space_before = Pt(4 if level == 0 else 2)
        run_b = p.add_run()
        run_b.text = "▸ " if level == 0 else "  · "
        run_b.font.size = Pt(12 if level == 0 else 10)
        run_b.font.bold = (level == 0)
        run_b.font.color.rgb = C_ACCENT if level == 0 else C_ACCENT2
        run_t = p.add_run()
        run_t.text = text
        run_t.font.size = Pt(12 if level == 0 else 10)
        run_t.font.color.rgb = C_WHITE if level == 0 else C_LIGHT_GREY

    _add_rect(slide, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), C_TITLE_BAR)
    _notes(slide, notes_text)
    return slide


def _summary_slide(prs, slide_id, title, rows, notes_text):
    """Table-style summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_SLIDE)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_ACCENT)
    _add_rect(slide, 0, Inches(0.06), SLIDE_W, Inches(0.72), C_BG_MID)

    tb_title = _add_textbox(slide, Inches(0.3), Inches(0.08), Inches(11.8), Inches(0.65))
    _para(_tf(tb_title), title, 20, True, C_WHITE, PP_ALIGN.LEFT)
    tb_id = _add_textbox(slide, Inches(12.2), Inches(0.1), Inches(1.0), Inches(0.5))
    _para(_tf(tb_id), slide_id, 8, False, C_ACCENT, PP_ALIGN.RIGHT)

    row_h = Inches(0.58)
    y0 = Inches(0.9)
    col_widths = [Inches(2.2), Inches(4.5), Inches(6.3)]
    col_xs = [Inches(0.2), Inches(2.5), Inches(7.1)]
    colours_cycle = [RGBColor(0x0F, 0x28, 0x3C), RGBColor(0x14, 0x2F, 0x44)]

    for i, row in enumerate(rows):
        bg = colours_cycle[i % 2]
        _add_rect(slide, col_xs[0], y0 + i * row_h, SLIDE_W - Inches(0.4), row_h, bg)
        for j, (cell, col_x, col_w) in enumerate(zip(row, col_xs, col_widths)):
            tb = _add_textbox(slide, col_x + Inches(0.1), y0 + i * row_h + Inches(0.05),
                              col_w - Inches(0.15), row_h - Inches(0.1))
            colour = C_ACCENT if j == 0 else (C_WHITE if j == 1 else C_LIGHT_GREY)
            size = 10 if j == 2 else 11
            _para(_tf(tb), str(cell), size, j == 0, colour, PP_ALIGN.LEFT)

    _add_rect(slide, 0, SLIDE_H - Inches(0.05), SLIDE_W, Inches(0.05), C_TITLE_BAR)
    _notes(slide, notes_text)
    return slide


def _closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, C_BG_DARK)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), C_ACCENT)
    _add_rect(slide, 0, SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), C_ACCENT)

    tb = _add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12.0), Inches(1.0))
    _para(_tf(tb), "Summary — Current State and Next Steps", 26, True, C_WHITE, PP_ALIGN.LEFT)

    _add_rect(slide, Inches(0.6), Inches(1.8), Inches(12.0), Inches(0.03), C_ACCENT2)

    items = [
        ("BUILT",    C_SUCCESS,  "Execution-reliable infrastructure (EXE-01 → EXE-05 certified)"),
        ("BUILT",    C_SUCCESS,  "Deterministic passive fill simulator with calibrated cost model"),
        ("BUILT",    C_SUCCESS,  "Forward validation pipeline with split/seed discipline (DAT-01 → VAL-03)"),
        ("BUILT",    C_SUCCESS,  "Pre-attempt gate mechanism — n_signals_before_gate diagnostic"),
        ("FOUND",    C_WARNING,  "Per-fill alpha exists marginally at fee = 0.5 bps"),
        ("FOUND",    C_DANGER,   "Per-attempt NPA negative at realistic fee = 1.0, adverse_mult = 1.0"),
        ("NEXT",     C_ACCENT,   "Gate calibration sweep: find NPA > 0 configuration on held-out splits"),
        ("GATE",     C_DANGER,   "Deployment blocked: NPA < 0  |  Trigger: NPA > 0 at fee=1.0, adv=1.0"),
    ]

    y = Inches(1.95)
    for status, colour, text in items:
        chip = _add_rect(slide, Inches(0.6), y, Inches(1.1), Inches(0.42), colour)
        tb_c = _add_textbox(slide, Inches(0.62), y + Inches(0.05), Inches(1.06), Inches(0.32))
        _para(_tf(tb_c), status, 9, True, C_BG_DARK, PP_ALIGN.CENTER)

        tb_t = _add_textbox(slide, Inches(1.85), y, Inches(11.2), Inches(0.42))
        _para(_tf(tb_t), text, 12, False, C_WHITE if colour != C_DANGER else RGBColor(0xFF, 0xAA, 0xAA),
              PP_ALIGN.LEFT)
        y += Inches(0.52)

    _notes(slide,
        "The infrastructure is production-ready. The research pipeline is scientifically rigorous. "
        "The alpha state is pre-deployment: NPA is negative at realistic conditions. "
        "When NPA turns positive on held-out data at realistic costs, the system is ready to deploy. "
        "Until then, the correct engineering answer is: measure more carefully, not deploy more aggressively."
    )
    return slide


# ---------------------------------------------------------------------------
# Full slide data — ALL sections
# ---------------------------------------------------------------------------

def build_all_slides(prs: Presentation):

    # ── TITLE ──────────────────────────────────────────────────────────────
    _title_slide(prs)

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 1 — SYSTEM ARCHITECTURE
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "I", "System Architecture",
                     "Eclipse Scalper as a Distributed Execution Organism")

    _content_slide(prs, "1.1", "Architectural Overview — The Execution Pipeline", [
        ("bootstrap.py — process lifecycle, lock acquisition, state hydration", 0),
        ("strategies/eclipse_scalper.py — per-symbol market event processing", 0),
        ("entry_loop.py — 20+ pre-submission gate checks before any order", 0),
        ("order_router.py — retry logic, idempotency, kill-switch enforcement", 0),
        ("exchanges/binance.py — ccxt wrapper, hedge-mode positionSide", 0),
        ("reconcile.py — exchange state as truth, local state as belief", 0),
        ("position_manager.py — stop-loss and take-profit lifecycle", 0),
        ("Every layer has a single responsibility; no downstream component can override safety", 1),
    ], notes_text=(
        "The architecture is deliberately linear in its happy path but defensively layered at every "
        "failure point. The signal generator produces raw directional signals; it has no authority to "
        "submit orders. Every signal must pass through entry_loop's gate complex before the router is "
        "called. The router is stateless — it acts on intents passed to it. Reconcile runs asynchronously "
        "and corrects divergence between local belief and exchange truth. The position manager handles "
        "only reductive orders and is exempt from entry restrictions by design.\n\n"
        "Kill-switch cannot be bypassed by any downstream component — it is checked independently at "
        "both entry_loop and router level."
    ), visual_label=(
        "Vertical pipeline: 6 boxes (bootstrap→signal→entry_loop→router→exchange→reconcile). "
        "Color: green=signal, amber=gates, blue=routing, grey=reconcile. "
        "Annotations show where kill-switch intercepts and where state reads occur."
    ))

    _content_slide(prs, "1.2", "The Brain State Model — Eventual Consistency Under Adversity", [
        ("PsycheState (brain/state.py): persisted as LZ4-compressed binary on disk", 0),
        ("Single source of truth: the exchange — local state is 'belief', not 'fact'", 0),
        ("run_context dict: ephemeral-but-durable runtime metadata surviving restarts", 0),
        ("Per-symbol async locks: shared between reconcile and position_manager", 0),
        ("WAL (write-ahead log) intents: prevent ghost submissions on restart", 0),
        ("Symbol canonicalization: BTC/USDT:USDT → BTCUSDT via symkey() — single source of truth", 0),
        ("Crash-safe without distributed transaction protocol — WAL + reconcile handles it", 1),
    ], notes_text=(
        "The state model adopts eventual consistency: local brain state converges to exchange truth "
        "on each reconciliation cycle. The WAL intent system prevents restart-induced double submissions. "
        "Before any order is submitted, a WAL intent is written to run_context. On restart, bootstrap "
        "reads open WAL intents and reconciles them before resuming. LZ4 compression keeps I/O bounded "
        "without meaningful latency."
    ), visual_label=(
        "Two-column: left=PsycheState (positions, run_context, WAL intents), "
        "right=Exchange state (fills, open orders). Bidirectional arrow labeled 'reconcile' between them. "
        "Crash scenario annotation: WAL intent survives → next boot resolves."
    ))

    _content_slide(prs, "1.3", "Kill-Switch and Circuit Breaker — Safety Hierarchy", [
        ("Kill-switch: binary halt — blocks ALL new entry submissions, always respected", 0),
        ("Circuit breaker: trip-on-loss — activates on threshold drawdown within rolling window", 0),
        ("EXITS ARE PRIVILEGED: intent_reduce_only=True orders bypass all entry restrictions", 0),
        ("No infinite retry loops: errors classified as retryable | fatal | idempotent_safe", 0),
        ("Kill-switch checked independently at both entry_loop AND router levels", 0),
        ("Circuit breaker is opt-in via bot.cfg.CIRCUIT_BREAKER_ENABLED — disabled in tests", 0),
        ("A kill-switch that also blocks stop-losses would amplify risk, not manage it", 1),
    ], notes_text=(
        "Kill-switch and circuit breaker are the top of the safety hierarchy. No component can override "
        "them for entry orders. This is enforced by architecture, not convention. The critical distinction: "
        "position reduction orders must always be allowed through. intent_reduce_only=True is the "
        "mechanism. Error classification prevents both missed fills (too strict) and order spam (too loose)."
    ), visual_label=(
        "Inverted triangle hierarchy: top (red)=Kill-Switch, middle (amber)=Circuit Breaker, "
        "bottom (green)=Normal Entry Flow. Separate lane on right: 'Exits/Reduce-Only' bypasses all levels. "
        "Annotations show the responsible module at each level."
    ))

    _content_slide(prs, "1.4", "Intent Lifecycle and Idempotency Guarantees", [
        ("Intent: typed record describing desired position action (open, reduce, close)", 0),
        ("States: pending → submitted → confirmed → completed | failed | orphaned", 0),
        ("Idempotency: same intent submitted twice → exchange detects via clientOrderId", 0),
        ("clientOrderId: deterministic, derived from symbol+side+timestamp — max 35 chars (Binance limit)", 0),
        ("Orphan detection: intents with no exchange match after TTL → reconciled away", 0),
        ("EXE-02 Lifecycle completeness: every intent must reach a terminal state — no leaks", 0),
        ("WAL intent written BEFORE exchange call — crash between write and submit is handled", 1),
    ], notes_text=(
        "The intent lifecycle is the core abstraction that makes the system restartable and observable. "
        "An intent is created before any exchange interaction. If the process crashes between intent "
        "creation and order submission, the WAL intent survives and bootstrap handles it on restart. "
        "clientOrderId is critical: Binance's idempotency guarantee. The 35-character limit shapes ID "
        "generation. EXE-02 prevents intent leaks — a class of bug where intents accumulate in pending "
        "state indefinitely, causing progressively incorrect position accounting."
    ), visual_label=(
        "State machine: nodes=pending/submitted/confirmed/completed/failed/orphaned. "
        "Directed edges labeled: submit, exchange ack, reconcile fill, reconcile miss TTL, exchange reject. "
        "Terminal states color-coded: completed=green, failed=amber, orphaned=red. "
        "Inset: clientOrderId format string."
    ))

    _content_slide(prs, "1.5", "Execution Invariants — The Non-Negotiable Contracts", [
        ("EXE-01 Router idempotency: same intent → same observable result, no side-effect doubling", 0),
        ("EXE-02 Lifecycle completeness: every intent reaches terminal state, no leaks ever", 0),
        ("EXE-03 Kill-switch precedence: kill-switch overrides all other logic for entries", 0),
        ("EXE-04 Restart safety: state after restart ≡ state of continuous operation", 0),
        ("EXE-05 Risk bounds: position sizes and exposure never exceed configured maximums", 0),
        ("A proposed change that violates any of these five is rejected regardless of claimed benefit", 1),
        ("EXE-04 is the hardest: run 10h + crash + restart ≡ run 20h continuously", 1),
    ], notes_text=(
        "These five invariants are the load-bearing walls of the execution system. Every architectural "
        "decision is evaluated against them. EXE-01 is hardest to test rigorously — requires reasoning "
        "across network partitions. WAL intents + bounded retries + clientOrderId idempotency approximate "
        "this. EXE-05 is enforced at entry_loop by checking current exposure against configured maximums "
        "after reconcile has run, reflecting exchange-confirmed state."
    ), visual_label=(
        "Five rows, each with invariant ID, name, one-line description, and green checkmark. "
        "Footer: 'Violation of any invariant = critical system failure requiring immediate halt'."
    ))

    _content_slide(prs, "1.6", "Why Execution Reliability Outweighs Signal Accuracy", [
        ("A 55% accurate signal with 3% execution error rate → negative expected value", 0),
        ("Silent execution bugs compound: ghost orders, missed closes, position drift", 0),
        ("Measurement fidelity determines whether alpha exists or is an artifact", 0),
        ("Signal alpha is fragile: market regimes shift. Execution bugs are permanent.", 0),
        ("The system is a measurement instrument first, a trading system second", 0),
        ("Example: 55% accuracy, +0.8 bps avg → 2% wrong-size errors + 1% market-order slips → net negative", 1),
        ("Before claiming alpha exists, prove the measurement is not contaminated", 1),
    ], notes_text=(
        "In a highly competitive market microstructure environment, edge from signal alpha is typically "
        "very small — measured in basis points per attempt. Any systematic execution error, even small, "
        "can entirely consume or invert that edge. The deeper implication: before claiming any alpha "
        "exists, prove that measurement is not contaminated by execution artifacts. This is why "
        "reconciliation-first architecture, deterministic simulation, and invariant checking exist — "
        "not as niceties but as epistemological necessities."
    ), visual_label=(
        "Two bar charts: left=Signal Accuracy Distribution (45/50/55/60%) with theoretical EV. "
        "Right=Execution Error Impact: how 0/1/2/3% error rates erode EV at each accuracy level. "
        "Red 'EV=0' line overlaid. Shows 55% signals go negative with modest execution error."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 2 — ALPHA DISCOVERY PIPELINE
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "II", "Microstructure Alpha Discovery Pipeline",
                     "From Raw Tick Data to Forward-Validated Signals")

    _content_slide(prs, "2.1", "The Research Pipeline Architecture", [
        ("Stage 0: data/microstructure.db — SQLite store of trade ticks and mark prices", 0),
        ("Stage 1: micro_edge_lib.py — bucket feature generation (1-second buckets)", 0),
        ("Stage 2: micro_edge_backtest.py — simulate_rule_trades with passive fill model", 0),
        ("Stage 3: sweep_*.py — parameter grid search: thresholds × fees × adverse multipliers", 0),
        ("Stage 4: validate_passive_pocket_forward.py — true forward splits, zero training contamination", 0),
        ("Stage 5: rank_passive_pockets_forward.py — capacity-adjusted NPA-based ranking", 0),
        ("Pipeline is strictly sequential; each stage has a defined API contract", 1),
        ("Database is the only shared persistent artifact; all intermediates are in-memory or reports", 1),
    ], notes_text=(
        "The 1-second bucket granularity reflects a design choice: sub-second resolution lets noise "
        "dominate signal. At 1-second buckets, we capture meaningful intensity and imbalance patterns "
        "while keeping feature dimensionality manageable. The passive fill simulator separates this "
        "pipeline from a naive backtester — it models fill probability given observable microstructure "
        "at signal time, accounting for queue position, trade intensity, and adverse selection."
    ), visual_label=(
        "Horizontal pipeline: 6 labeled stages with arrows. Each shows module name, input type, "
        "output type. Below each stage: data icon (DB cylinder, RAM chip, document). "
        "Arrows labeled with interface: 'bucket rows', 'sim results', 'eval grid', 'ranked list'."
    ))

    _content_slide(prs, "2.2", "Microstructure Features — What the System Observes", [
        ("Trade Intensity: Σ(trade_count × volume) / bucket_sec — informed flow arrival rate proxy", 0),
        ("Order Book Imbalance: (bid_depth − ask_depth) / (bid_depth + ask_depth) — directional pressure", 0),
        ("Spread Compression: spread / mid_price — queue eligibility and fill probability proxy", 0),
        ("Micro-Volatility: rolling std-dev of 1-bucket returns — regime indicator", 0),
        ("V2 Suite: persistence-weighted imbalance score, mean-reversion probability, flip rate", 0),
        ("High-volatility regimes → toxicity gate activates → passive orders blocked (adverse selection)", 1),
        ("V2 imbalance persistence: consecutive buckets above threshold → stronger signal", 1),
    ], notes_text=(
        "Trade intensity is the primary signal driver. High-intensity periods indicate elevated informed "
        "order flow arrival. The hypothesis: sufficiently intense + directionally imbalanced flow predicts "
        "short-term price continuation, and a passive maker can profitably sit on the predicted side. "
        "Spread compression is critical for passive viability: when spread is tight, passive fills are "
        "more likely because crossed trades frequently touch best bid/ask. Wide spreads reduce fill "
        "probability and increase adverse selection risk."
    ), visual_label=(
        "2×3 grid of mini time-series charts. Each cell: feature name, formula, representative 1h plot. "
        "Highlight: imbalance spike → price move; intensity spike → volatility; spread compression → fills."
    ))

    _content_slide(prs, "2.3", "The Passive Execution Simulator — Modeling Fill Reality", [
        ("Core: simulate_passive_fill(event, horizon_sec, features, params) in passive_execution_simulator.py", 0),
        ("Fill decision: u01 = _u01_from_key(seed, event_id) — deterministic uniform hash, not PRNG", 0),
        ("Touch probability: base × feature-conditional multipliers (spread, intensity, imbalance, vol)", 0),
        ("Queue competition: p_touch_adj = p_touch × (1 − queue_competition_strength × intensity_factor)", 0),
        ("Adverse selection: adverse_bps from calibrated conditional distribution — applied on fills ONLY", 0),
        ("Cost model: net_return = raw_return − (2 × maker_fee_bps/10000) − (adverse_bps/10000)", 0),
        ("Unfilled attempts: net_return = 0.0 exactly — not a cost, not a gain", 1),
    ], notes_text=(
        "The simulator is the epistemological core of the research pipeline. It answers: 'Given that my "
        "signal fired at time t, would a passive limit order have filled, and at what effective cost?' "
        "The _u01_from_key hash maps each (seed, event_id) pair independently to [0,1] — fill decision "
        "for event E42 is independent of whether events E1-E41 were simulated. This allows partial re-runs "
        "and event-level analysis without disturbing global random state."
    ), visual_label=(
        "Flow diagram: simulate_passive_fill internals. Input→p_touch calc→queue_competition adj→"
        "u01 draw→fill decision branch. Filled→adverse_bps sampled→cost applied→net_return. "
        "Unfilled→net_return=0.0. Different colors for filled/unfilled outcome boxes."
    ))

    _content_slide(prs, "2.4", "Deterministic Simulation — Reproducibility as Epistemology", [
        ("All randomness: _u01_from_key(seed, event_id) — pure hash function, no PRNG state", 0),
        ("Same (seed, event sequence) → identical fill decisions, costs, and returns — always", 0),
        ("Seeds are explicit parameters, never sourced from wall-clock or OS entropy", 0),
        ("attempts_per_min = n_attempts / (val_rows × bucket_sec / 60) — zero datetime.now() calls", 0),
        ("DAT-03 invariant: deterministic per (seed, input) — verified by test_net_per_attempt_deterministic", 0),
        ("DAT-01 invariant: no lookahead — features computed only from data ≤ signal time", 0),
        ("Without determinism: cannot distinguish genuine alpha from lucky seeds", 1),
    ], notes_text=(
        "Reproducibility is not a convenience feature; it is an epistemological requirement. Without "
        "determinism, we cannot distinguish genuine alpha from lucky random seeds. Without determinism, "
        "tests cannot catch regressions. The hash design (vs sequential PRNG) means partial re-runs "
        "and event-level analysis are possible without disturbing global state. DAT-01 is verified "
        "structurally: the feature computation function receives only rows up to the current signal index."
    ), visual_label=(
        "Side-by-side comparison: left='Non-Deterministic Approach' (PRNG seeded with wall-clock, "
        "different results each run). Right='Eclipse Scalper Approach' (hash-based, identical for "
        "identical inputs). Below: timing diagram showing feature inputs strictly at t≤signal_time."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 3 — PER-FILL ILLUSION
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "III", "The Per-Fill Illusion Problem",
                     "Why Conditional Profitability Is Not Expected Value")

    _content_slide(prs, "3.1", "Two Metrics, One Question: Do We Have Alpha?", [
        ("filled_avg_net = Σ(net_return | filled=True) / n_fills — conditional on execution", 0),
        ("net_per_attempt = Σ(net_return for ALL attempts) / n_attempts — unconditional EV", 0),
        ("Unfilled attempts contribute exactly 0.0 net return — not a loss, not a gain", 0),
        ("Mathematical relationship: NPA = fill_rate × filled_avg_net", 0),
        ("These metrics diverge when fill_rate < 1.0 — which is always true for passive orders", 0),
        ("Passive orders only fill when market moves to them — by definition fill_rate < 1", 1),
        ("filled_avg_net is a biased estimator of strategy profitability; NPA is unbiased", 1),
    ], notes_text=(
        "The distinction between filled_avg_net and net_per_attempt is the single most important "
        "conceptual point in this presentation. It is the difference between 'do we profit when we "
        "trade?' and 'do we profit in expectation across all trading decisions?'. A passive limit order "
        "strategy has fill_rate < 1.0 by definition. filled_avg_net is a conditional expectation: "
        "E[return | fill]. NPA = E[return] across all attempts. The formula NPA = fill_rate × "
        "filled_avg_net is exact. When fill_rate = 0.5 and filled_avg_net = +0.8 bps, NPA = +0.4 bps "
        "— but after costs this can go negative."
    ), visual_label=(
        "Funnel diagram: 100 signal events → 50 fills (fill_rate=50%). filled_avg_net computed from 50. "
        "Same 100 events with unfilled contributing 0.0 → NPA from all 100. "
        "Bar chart: filled_avg_net=+0.8bps vs NPA=-0.2bps after costs. The sign flip is the key message."
    ))

    _formula_slide(prs, "3.2", "Mathematical Decomposition of Strategy Expected Value",
        formulas=[
            "# Definitions",
            "N  = total attempts",
            "F  = fills",
            "f  = F / N   (fill rate)",
            "",
            "filled_avg_net",
            "  = (1/F) * Σ r_i  [i: filled]",
            "",
            "NPA = (1/N) * Σ r_i  [all i]",
            "    = f * filled_avg_net",
            "",
            "# Cost per fill (round trip)",
            "c_fee = 2 * maker_fee_bps / 10_000",
            "c_adv = adverse_bps / 10_000",
            "r_i   = raw_ret_i - c_fee - c_adv",
            "",
            "# Break-even fill rate",
            "f* = (c_fee + c_adv) / raw_filled_avg",
        ],
        explanation_bullets=[
            ("NPA is the unbiased estimator of strategy profitability", 0),
            ("filled_avg_net is biased — it conditions on fills occurring", 0),
            ("Break-even fee: filled_avg_net must exceed c_fee + c_adv", 0),
            ("At maker_fee=1.0 bps, c_fee=0.0001 (round trip)", 0),
            ("Adverse ~0.1 bps → c_adv=0.00001", 0),
            ("If raw_filled_avg=0.8 bps < 1.1 bps cost → every fill loses money", 0),
            ("More fills = more losses in this regime", 1),
            ("The gate is the tool to invert this", 1),
        ],
        notes_text=(
            "The break-even fill rate formula shows that strategy viability depends on the ratio of "
            "costs to gross return per fill. Our empirical estimates: maker fee = 1.0 bps round trip, "
            "adverse = 0.1 bps, total cost ≈ 1.1 bps. If raw average gross return per fill is 0.8 bps, "
            "we need gross returns > 1.1 bps to cover costs. Every fill currently loses money. "
            "NPA = fill_rate × (−0.3 bps) ≈ −0.15 bps. The gate's goal: select only signals with "
            "expected gross return > 1.1 bps."
        ))

    _content_slide(prs, "3.3", "Adverse Selection — The Hidden Cost of Passive Execution", [
        ("Adverse selection: we fill disproportionately when market moves against us immediately after", 0),
        ("Mechanism: informed traders hit our passive orders at moments of maximum adverse information", 0),
        ("adverse_bps calibrated from historical touch events — not assumed, not constant", 0),
        ("passive_adverse_mult parameter scales severity — 1.2 = +20% stress condition", 0),
        ("Adverse selection cost applied ONLY on fills — verified by test_npa_cost_applied_to_fills_only", 0),
        ("Adverse cost test: assert all unfilled attempts have net_return == 0.0 exactly", 1),
        ("test_adverse_bps_scaling_not_amplified: guards against 10,000× unit error in cost application", 1),
    ], notes_text=(
        "Adverse selection is the most subtle and damaging cost in passive execution strategies. It arises "
        "from asymmetric information: when an informed trader has directional information, they trade "
        "aggressively. Our passive order at the bid gets hit precisely when the informed trader has "
        "determined price should fall. We fill, and immediately price moves against us. The calibration "
        "is empirical: historical events where a passive order would have been touched → distribution of "
        "adverse moves in subsequent horizon window."
    ), visual_label=(
        "Two panels: left=time series showing price drop after adverse fill, annotated 'adverse_move=0.1bps'. "
        "Right=histogram of adverse_bps across 1000 simulated events, right-skewed, "
        "with calibrated mean and 90th percentile marked."
    ))

    _content_slide(prs, "3.4", "Attempt Spam — Why More Signals Is Not Better", [
        ("High signal frequency → more attempts → NPA dominated by fill costs, not fill returns", 0),
        ("If fill_rate=0.5 and all fills lose money → doubling signals doubles losses", 0),
        ("NPA is invariant to signal frequency ONLY if fills are profitable", 0),
        ("attempts_per_min metric quantifies signal spam rate — high values are warning signs", 0),
        ("Gate solution: reduce attempts to only those with high fill-quality expectation", 0),
        ("Every attempt in loss regime: E[return] = fill_rate × negative_filled_avg_net < 0", 1),
        ("200 attempts/min at NPA=−0.0001 → systematic and measurable loss rate", 1),
    ], notes_text=(
        "The attempt spam problem arises when a rule fires frequently but only a subset of firings "
        "correspond to genuine high-quality microstructure signals. Each firing becomes an attempt. "
        "Even unfilled attempts have zero return, not negative return. But filled attempts in low-quality "
        "conditions have negative return (costs exceed gross return). So E[value of attempting in "
        "low-quality conditions] = fill_rate × negative_return < 0. More attempts makes this worse. "
        "The gate is the solution: quality filter that increases average attempt quality."
    ), visual_label=(
        "Two time-series plots, same rule, same window. Top: without gate — 200 attempts/min, "
        "dense marks on price chart, cumulative PnL = steady decline. "
        "Bottom: with gate — 40 attempts/min, same quality fills remain, cumulative PnL = flat or positive."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 4 — RANKING EVOLUTION
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "IV", "Ranking Evolution",
                     "From Degenerate Scores to Capacity-Honest Evaluation")

    _content_slide(prs, "4.1", "Original Ranking — Three Systematic Flaws", [
        ("Flaw 1 — Conditional metrics only: filled_avg_net ignores fill rate entirely", 0),
        ("Flaw 2 — Degeneracy under gating: score=0 when gate fails → all failing pockets identical", 0),
        ("Flaw 3 — No capacity exposure: attempts_per_min, val_attempts, fill rate never surfaced", 0),
        ("Result: ranking cannot distinguish 'barely failing' from 'catastrophically failing'", 0),
        ("Result: pockets with positive per-fill but negative NPA ranked above viable alternatives", 0),
        ("Original formula: pass_rate × median_bps × worst_split_bps", 1),
        ("When pass_rate below threshold → score=0 regardless of actual NPA magnitude", 1),
    ], notes_text=(
        "The original ranking formula: score = pass_rate × median_bps × worst_split_bps. When pass_rate "
        "fell below the robustness gate threshold, score was clamped to zero. Consequence: all pockets "
        "that failed the gate received score=0, regardless of how close they were to passing. This "
        "degeneracy meant ranking was uninformative precisely when most needed — when trying to identify "
        "which failing pockets were closest to viability. A pocket with NPA=−0.1 bps ranked identically "
        "to NPA=−50 bps."
    ), visual_label=(
        "Before/after comparison. Left: table with score column — many entries at score=0.00, "
        "different fill metrics, demonstrating degeneracy. Right: table with score+score_raw_core+"
        "score_raw_stress — zero-scored pockets now show visible differentiation. "
        "Red bars highlight degeneracy region in left table."
    ))

    _formula_slide(prs, "4.2", "The New Scoring Framework — NPA-First, Capacity-Aware",
        formulas=[
            "# Grid points",
            "fee_one  = fee closest to 1.0 bps",
            "adv_one  = adverse closest to 1.0",
            "fee_min  = min(fee_grid)",
            "fee_max  = max(fee_grid)",
            "adv_max  = max(adverse_grid)",
            "",
            "# Core & stress evals",
            "core_npa   = eval(fee_one, adv_one).median_NPA",
            "stress_npa = eval(fee_one, adv_max).median_NPA",
            "",
            "# Score computation",
            "npa_bps   = core_npa * 10_000",
            "stab_bps  = stability_std * 10_000",
            "base      = max(0, npa_bps) / (1 + max(0, stab_bps))",
            "cap_pen   = 1 + 0.5 * insufficient_fill_rate",
            "",
            "score = base / cap_pen",
            "     if (core_npa > 0 AND stress_npa > 0)",
            "     else 0",
        ],
        explanation_bullets=[
            ("NPA > 0 at core AND stress = robustness gate", 0),
            ("Higher NPA → higher base score", 0),
            ("Higher stability_std → lower score (variance penalized)", 0),
            ("Higher insufficient_fill_rate → capacity penalty", 0),
            ("score=0 only when NPA fails, not arbitrary threshold", 0),
            ("score_raw_* always populated — diagnose zero-scored pockets", 0),
            ("Sort key: (score DESC, score_raw_core DESC)", 1),
            ("Zero-gated pockets remain comparable via score_raw_core", 1),
        ],
        notes_text=(
            "The scoring formula embodies three principles. First, NPA is primary — pocket must have "
            "positive NPA at realistic conditions for non-zero score. Second, stability matters — "
            "high variance across splits is penalized. Third, capacity matters — pockets where high "
            "proportion of attempts fail to fill receive capacity penalty. The robustness gate "
            "(core_npa > 0 AND stress_npa > 0) replaced the old pass_rate gate. Now directly tied "
            "to the metric we care about."
        ))

    _content_slide(prs, "4.3", "Score Raw Fields — Diagnostic Transparency", [
        ("score_raw_core: median_filled_avg_net at (fee_min, adv≈1.0) — most favorable conditions", 0),
        ("score_raw_stress: median_filled_avg_net at (fee_max, adv_max) — worst-case conditions", 0),
        ("score_raw_min: min(median_filled_avg_net) across entire evaluation grid", 0),
        ("All three fields always populated — even when score = 0", 0),
        ("score_raw_core > 0, score_raw_stress < 0 → fragile: profitable but adverse-sensitive", 0),
        ("score_raw_core < 0 → fundamentally broken: unprofitable even under favorable conditions", 0),
        ("Different failure modes require different remediation strategies", 1),
    ], notes_text=(
        "The score_raw fields make the zero-score region informative. Before this change, 40 pockets "
        "could all have score=0 with no way to determine which were closest to viability. "
        "score_raw_core provides ordering within the zero-score cohort. A pocket with score_raw_core > 0 "
        "but score_raw_stress < 0: profitable under normal conditions but fragile under adverse selection. "
        "This is a different failure mode from score_raw_core < 0 (unprofitable even at favorable fee) "
        "and requires different remediation."
    ), visual_label=(
        "3×3 matrix heatmap: rows=adverse_mult (0.8,1.0,1.2), cols=fee (0.5,1.0,1.5 bps). "
        "Each cell=color-coded median_filled_avg_net (green=positive, red=negative). "
        "Circles mark: score_raw_core (fee_min, adv=1.0), score_raw_stress (fee_max, adv_max), "
        "score_raw_min (worst cell). Side-by-side ETHUSDT vs BTCUSDT heatmaps."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 5 — PRE-ATTEMPT GATE
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "V", "Pre-Attempt Gate Implementation",
                     "Signal Generation vs Execution Eligibility")

    _content_slide(prs, "5.1", "The Conceptual Separation: Signal vs Attempt", [
        ("Signal: rule fires — microstructure conditions met, direction predicted", 0),
        ("Attempt: passive limit order would be placed — eligibility check passed", 0),
        ("Signals are necessary but not sufficient for attempts", 0),
        ("Gate = filter between signal generation and attempt registration", 0),
        ("Gate operates on signal-time features only — DAT-01 compliant, zero lookahead", 0),
        ("Gate is deterministic per (row, gate_config) — same input → same decision", 0),
        ("Gate-blocked signals do NOT consume cooldown — next signal can fire immediately", 1),
        ("Implemented in _passes_attempt_gate(r, gate) in tools/micro_edge_backtest.py", 1),
    ], notes_text=(
        "The gate introduces a second-level filter between rule firing and attempt registration. "
        "A signal fires when the rule's primary conditions are met. An attempt is registered only when "
        "the gate's additional strength conditions are also met. DAT-01 compliance is structural: "
        "_passes_attempt_gate(r, gate) receives only the current bucket row r, computed from data "
        "available at signal time. The gate's non-interaction with cooldown is a critical design choice: "
        "blocked signals don't consume the cooldown window, keeping the system responsive to the "
        "next high-quality signal."
    ), visual_label=(
        "Two-tier flowchart. Tier 1 (Signal Layer): rule_fires→feature_bounds_check→trade_side_resolved. "
        "Tier 2 (Gate Layer): _passes_attempt_gate with three decision diamonds (intensity/imbalance/spread). "
        "PASS → 'Attempt Registered' (green). BLOCK → counters incremented, advance i (red). "
        "Horizontal dashed line labeled 'Signal/Attempt Boundary'."
    ))

    _content_slide(prs, "5.2", "Gate Criteria — Selecting for Execution Quality", [
        ("min_trade_intensity_strong: attempt only when trade_intensity > threshold", 0),
        ("min_imbalance_strong: attempt only when |imbalance| > threshold", 0),
        ("max_spread_tight: attempt only when spread < threshold", 0),
        ("All criteria use 'or 0.0' fallback for None feature values — safe, not strict-fail", 0),
        ("Default: all zero — gate disabled, fully backward-compatible behavior", 0),
        ("Gate is additive: all non-zero criteria must pass (AND logic)", 0),
        ("High-intensity: order book turns over faster → passive fills at better prices", 1),
        ("Tight-spread: bid-ask bounce cost scales with spread → lower adverse selection", 1),
    ], notes_text=(
        "The three gate criteria correlate with execution quality in passive limit order strategies. "
        "min_trade_intensity_strong: during elevated intensity, passive orders are more likely to fill "
        "at favorable prices before market moves decisively. min_imbalance_strong: higher imbalance "
        "magnitude indicates stronger directional pressure — more predictive of short-term continuation. "
        "max_spread_tight: narrow spread → smaller adverse selection component relative to gross return. "
        "AND logic means we want the INTERSECTION of favorable conditions, not any one in isolation."
    ), visual_label=(
        "Three horizontal sliders with distribution curves above each. "
        "Slider 1: trade_intensity distribution, threshold at 'strong' level, shaded='passes'. "
        "Slider 2: |imbalance| distribution with threshold. "
        "Slider 3: spread distribution (inverted, tight=right). "
        "Below: Venn diagram with three circles; intersection=green 'attempt region'; rest=blocked."
    ))

    _content_slide(prs, "5.3", "Gate Impact Metrics — Before/After Measurement", [
        ("n_signals_before_gate: signals reaching gate — total pre-filter count", 0),
        ("n_attempts / val_attempts_after_gate: count of attempts post-gate", 0),
        ("gate_pass_rate = n_attempts / n_signals_before_gate — selectivity measure", 0),
        ("attempt_gate_blocked = n_signals_before_gate − n_attempts — blocked count in debug_stats", 0),
        ("val_attempts_before_gate vs val_attempts_after_gate: visible per combo in reports", 0),
        ("net_per_attempt_after_gate: NPA on gated attempt pool — the target metric", 0),
        ("Goal: find gate config where net_per_attempt_after_gate > 0 on held-out splits", 1),
        ("Target gate_pass_rate: 20–60% — over-restrictive (<5%) or under-filtering (>95%) are both wrong", 1),
    ], notes_text=(
        "The diagnostic metrics are essential for gate calibration. Without them, we cannot know whether "
        "a gate configuration is appropriately selective or over-aggressive. gate_pass_rate is the "
        "primary calibration indicator. net_per_attempt_after_gate determines whether the gate actually "
        "improves EV — this is testable: run with and without gate on the same validation split, "
        "measure NPA in both cases. test_attempt_gate_reduces_n_attempts verifies structural correctness; "
        "whether the gate improves NPA is an empirical question for the forward validation pipeline."
    ), visual_label=(
        "Two-panel bar chart. Left: bars per seed/split showing n_signals total height, "
        "split into 'passed gate' (green) and 'blocked' (red). "
        "Right: NPA per seed/split — with gate (green) vs without gate (amber). "
        "Expected result when calibration works: green bars are higher (less negative or positive)."
    ))

    _code_slide(prs, "5.4", "Gate Implementation — Core Code Paths", [
        "# tools/micro_edge_backtest.py",
        "",
        "def _passes_attempt_gate(r: Dict[str, Any], gate: Dict[str, float]) -> bool:",
        "    '''Signal-time features only. DAT-01 safe. No lookahead.'''",
        "    if not gate:",
        "        return True",
        "    ti  = float(r.get('trade_intensity') or 0.0)",
        "    imb = abs(float(r.get('imbalance') or 0.0))",
        "    spr = float(r.get('spread') or 0.0)",
        "    if gate.get('min_trade_intensity_strong', 0.0) > 0.0:",
        "        if ti < gate['min_trade_intensity_strong']: return False",
        "    if gate.get('min_imbalance_strong', 0.0) > 0.0:",
        "        if imb < gate['min_imbalance_strong']: return False",
        "    if gate.get('max_spread_tight', 0.0) > 0.0:",
        "        if spr > gate['max_spread_tight']: return False",
        "    return True",
        "",
        "# In simulate_rule_trades — insertion point after trade_side resolved:",
        "n_pre_gate_signals += 1",
        "if attempt_gate_bounds and not _passes_attempt_gate(r, attempt_gate_bounds):",
        "    debug_stats['attempt_gate_blocked'] += 1",
        "    i += 1",
        "    continue   # no cooldown consumed — gate block is transparent to scheduler",
        "",
        "# attempt_level_metrics now includes:",
        "# 'n_signals_before_gate': int(n_pre_gate_signals)",
    ], notes_text=(
        "The gate function is a pure function with no side effects. It reads only the current row r "
        "and the gate configuration dict. The insertion point in simulate_rule_trades is critical: "
        "after trade_side is resolved (so we know direction) but before the attempted={} dict is created "
        "(so gate-blocked signals never enter attempt_rows). The continue without updating next_allowed "
        "means gate-blocked signals do not consume the cooldown window."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 6 — FEE AND ADVERSE SENSITIVITY
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "VI", "Fee and Adverse Sensitivity Analysis",
                     "Mapping the Viability Surface")

    _content_slide(prs, "6.1", "The Fee-Adverse Grid — Stress-Testing Alpha Robustness", [
        ("Grid axes: maker_fee_bps ∈ {0.5, 1.0, 1.5} × passive_adverse_mult ∈ {0.8, 1.0, 1.2}", 0),
        ("9 evaluation points per pocket — full cross-product, same seeds and splits", 0),
        ("Core: (fee=1.0, adv=1.0) — realistic baseline for ranking and NPA gate", 0),
        ("Stress: (fee=1.0, adv=1.2) — +20% adverse selection severity", 0),
        ("Extreme: (fee=1.5, adv=1.2) — worst-case: maximum fee + maximum adverse", 0),
        ("Maker fee range: covers Binance fee tier uncertainty and rebate program changes", 1),
        ("Adverse mult range: covers calibration uncertainty ±20% — structural sensitivity check", 1),
    ], notes_text=(
        "The fee-adverse grid is a systematic stress test of strategy alpha across two independent cost "
        "dimensions. The maker fee dimension captures exchange fee structure uncertainty. Binance rebates "
        "vary by volume tier and can shift the effective fee substantially. The adverse multiplier "
        "dimension captures calibration uncertainty. The 9-point grid allows mapping the fee-adverse "
        "surface and identifying breakpoints. A robust pocket shows positive NPA across most of the grid. "
        "A fragile pocket shows positive NPA only at the (0.5, 0.8) corner — not deployable."
    ), visual_label=(
        "3×3 heatmap: fee on x-axis (0.5, 1.0, 1.5), adverse_mult on y-axis (0.8, 1.0, 1.2). "
        "Each cell color-coded by median_net_per_attempt: green=positive, red=negative. "
        "Black border around core cell (1.0, 1.0). "
        "Two heatmaps side-by-side: ETHUSDT and BTCUSDT for comparison."
    ))

    _content_slide(prs, "6.2", "Break-Even Fee Analysis — Where Alpha Inverts", [
        ("Break-even fee: maker_fee where NPA crosses zero (adverse held constant)", 0),
        ("best_fee_survive: max fee where pass_rate ≥ 50% — surfaced in ranking output", 0),
        ("Fragile alpha: NPA positive only at fee=0.5 bps — fee-tier arbitrage, not structural edge", 0),
        ("Robust alpha: NPA positive at fee=1.0 bps and beyond — genuine microstructure signal", 0),
        ("Current empirical finding: best_fee_survive=0.0 for all pockets — NPA < 0 at all fees", 0),
        ("Fee-dependent strategy: viable only with maker rebates — program-change risk", 1),
        ("Deployment requirement: NPA > 0 at fee=1.0, adv=1.0 — no rebate dependency", 1),
    ], notes_text=(
        "The break-even fee concept is central to assessing whether observed alpha is a genuine "
        "structural edge or a fee-tier artifact. Some passive strategies appear profitable only at "
        "unrealistically low fee assumptions. best_fee_survive: highest fee in the grid at which "
        "pocket still achieves pass_rate >= 50%. Currently all pockets show best_fee_survive=0.0 — "
        "consistent with negative NPA finding. Practical implication: a strategy requiring maker "
        "rebates to be profitable is not truly a market-making alpha — it is rebate harvesting with "
        "different risks."
    ), visual_label=(
        "Line chart: x=maker_fee_bps (0.5 to 1.5), y=NPA (×10⁻⁴). "
        "Multiple lines, one per pocket configuration. "
        "Horizontal zero line. x-intercept of each line = break-even fee. "
        "Green lines: cross zero at fee>1.0 (robust). Red lines: cross zero at fee<0.5 (fragile). "
        "Label break-even crossing points."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 7 — DETERMINISM AND SCIENTIFIC VALIDITY
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "VII", "Determinism and Scientific Validity",
                     "Engineering Reproducibility as Quantitative Method")

    _summary_slide(prs, "7.1", "Invariant Taxonomy — Data and Research Contracts", [
        ("DAT-01", "No Lookahead",          "Features computed only from data ≤ signal time — verified structurally"),
        ("DAT-02", "Timing Alignment",       "Feature timestamps aligned to bucket close, not trade arrival"),
        ("DAT-03", "Deterministic Sim",      "Same (seed, input) → same fill decisions, NPA — test_net_per_attempt_deterministic"),
        ("DAT-04", "Cost Unit Correctness",  "adverse_bps in bps units, applied as ratio/10000 — test_adverse_bps_scaling"),
        ("DAT-05", "Schema Stability",       "Additive-only JSONL changes: no renames, no removals ever"),
        ("VAL-01", "True Forward Splits",    "Validation window never overlaps training window — enforced by split_ranges()"),
        ("VAL-02", "Ranking Reproducibility","Same seed → same NPA → same rank — verified by determinism tests"),
        ("VAL-03", "Candidate Parsing",      "Parser handles alias columns, deduplication, YES/true/1 pass values"),
        ("SAF-01", "Exit Privilege",         "reduce_only orders bypass all entry gates — never blocked by kill-switch"),
        ("SAF-02", "Bounded Retries",        "No infinite retry loops — errors classified retryable/fatal/idempotent"),
    ], notes_text=(
        "The data and research invariants are the scientific methodology encoded in software. "
        "DAT-01 is the most critical and most commonly violated invariant in backtesting systems. "
        "Lookahead bias occurs when a model has access to information at time t unavailable to a live "
        "trader. In our system, this is prevented structurally. DAT-03 is verified by test: two calls "
        "with identical inputs and seeds must produce identical per_combo outputs. DAT-04 is subtle but "
        "financially critical: if adverse_bps=0.1 is applied as 0.1 (not 0.1/10000), cost model "
        "overstates adverse selection by 10,000× — making every fill appear catastrophically unprofitable."
    ))

    _content_slide(prs, "7.2", "Test Architecture — 79 Executable Specifications", [
        ("79 tests in tests/ directory — all pass as of 2026-02-20", 0),
        ("Framework: pytest with monkeypatching for full dependency isolation", 0),
        ("Isolation pattern: patch _load_symbol_trades_and_marks, build_bucket_features, simulate_rule_trades", 0),
        ("Tests are deterministic: fixed seeds, fixed mock outputs, no time.time() in hot paths", 0),
        ("New tests added this session: test_npa_cost_applied_to_fills_only, test_attempt_gate_reduces_n_attempts", 0),
        ("test_npa_cost_applied_to_fills_only: assert all unfilled attempts have net_return==0.0 exactly", 1),
        ("test_attempt_gate_reduces_n_attempts: assert n_gated < n_no_gate and gate_blocked > 0", 1),
    ], notes_text=(
        "The test suite serves dual purposes: regression prevention and specification. A test is an "
        "executable specification of expected behavior. The monkeypatching strategy is particularly "
        "important: tests cannot require a live database. Every database-reading function is patched "
        "with deterministic synthetic data. The 79 tests grew from 77 (added 2 this session). "
        "Count growth tracks feature additions: every new capability must have at least one test. "
        "When 79 tests pass, we have 79 executable certifications of system properties."
    ), visual_label=(
        "Test coverage matrix. Rows=components (cost_model/fill_sim/validator/ranker/gate). "
        "Columns=invariant types (DAT-01 through VAL-03, EXE-01 through SAF-02). "
        "Cells: filled=tested, empty=not tested, half=partial. "
        "Test count annotation per row. Footer: '79 passing — 0 failing'."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 8 — CURRENT FINDINGS
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "VIII", "Current Findings and Reality Check",
                     "Honest Assessment of Present Alpha State")

    _content_slide(prs, "8.1", "Empirical Results — What the Data Actually Shows", [
        ("Per-fill: filled_avg_net ≈ +2×10⁻⁵ to +4×10⁻⁵ at fee=0.5, adv=1.0 — marginal positive", 0),
        ("Per-attempt: net_per_attempt ≈ −6×10⁻⁵ to −1×10⁻⁴ at realistic conditions — consistently negative", 0),
        ("Attempt fill rate: ~50% — half of passive order attempts result in fills", 0),
        ("score = 0 for all ranked pockets: robustness gate fails at core conditions", 0),
        ("score_raw_core > 0 for top pockets: per-fill edge exists at favorable fee conditions", 0),
        ("score_raw_stress < 0 for all pockets: adverse selection stress inverts the edge", 0),
        ("Interpretation: costs (fee + adverse) dominate per-fill gross returns", 1),
        ("filled_avg_net ≈ 0.4 bps < 1.1 bps total cost → net fill return is negative", 1),
    ], notes_text=(
        "These numbers represent the unvarnished truth. The per-fill edge (filled_avg_net > 0) exists "
        "— the signal is directionally informative and passive orders do, on average, profit slightly "
        "relative to zero BEFORE costs. The problem: this per-fill edge is not large enough to survive "
        "full cost structure. At fee=0.5 bps, round-trip fee = 1.0 bps. If filled_avg_net ≈ 0.3 bps "
        "(before costs), net filled return = 0.3 − 1.0 − 0.1 = −0.8 bps. NPA = 0.5 × (−0.8 bps) = "
        "−0.4 bps ≈ −4×10⁻⁵. This matches the empirical NPA range."
    ), visual_label=(
        "Two-panel bar chart. Left: per-fill results by pocket — some positive bars at fee=0.5 (green). "
        "Right: per-attempt results at fee=1.0 — all bars clearly negative (red). "
        "Zero line on both panels. Annotations: 'fill_rate=50%', 'fee_cost=1.0bps'. "
        "Arrow labeled 'NPA = fill_rate × (filled_avg_net − costs)' showing the collapse."
    ))

    _two_col_slide(prs, "8.2",
        "Why Negative NPA Is a System Success, Not a System Failure",
        left_items=[
            "System correctly measured an illusion — did not deploy it",
            "Alternative: deploy without NPA awareness → guaranteed live losses",
            "Measurement fidelity is the primary product",
            "False positive deployment rate = 0%",
            "Pipeline's value: what it prevents as much as what it finds",
            "Every negative result narrows the search space for genuine alpha",
        ],
        right_items=[
            "Counterfactual: deploy on per-fill metrics → live fee losses + adverse",
            "Exchange fees are real and immediate — no simulation artifact",
            "The pipeline identified failure mechanism precisely: cost domination",
            "Quantified the gap: NPA ≈ −0.1 bps at realistic conditions",
            "Proposed the remediation: pre-attempt gate to raise quality bar",
            "Next hypothesis is falsifiable: gate → NPA > 0 on held-out splits",
        ],
        notes_text=(
            "This reframing is important. A naive interpretation of 'we found no profitable strategy' "
            "is that the system failed. The correct interpretation: the system successfully detected "
            "that what appeared to be alpha (positive per-fill metrics) was NOT genuine alpha "
            "(negative NPA). This is exactly what the pipeline is designed to do. The scientific "
            "process: hypothesis=signal generates alpha, experiment=forward validation with realistic "
            "costs, result=negative, conclusion=signal as configured doesn't overcome costs, "
            "next hypothesis=gate raises quality enough for positive NPA. That is testable."
        ),
        left_header="System Achievements",
        right_header="What Comes Next")

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 9 — NEXT EVOLUTION PATH
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "IX", "Next Evolution Path",
                     "Engineering the Path to Positive NPA")

    _content_slide(prs, "9.1", "Signal Refinement — Raising Selection Quality", [
        ("Current signal: intensity_spike_imbalance_cont — broad threshold rule, high firing rate", 0),
        ("V2 features: v2_min_score, v2_min_persistence, v2_min_confidence — richer information space", 0),
        ("v2_imbalance_persist: consecutive buckets above threshold → low-noise signal filter", 0),
        ("v2_flip_rate: rapidly reversing imbalance → unreliable signal; low flip_rate = stable", 0),
        ("Hypothesis: signals scoring highly on V2 features have materially higher filled_avg_net", 0),
        ("Test: compare filled_avg_net distribution at V2 thresholds vs base thresholds on same window", 1),
        ("If confirmed: V2 filtering converts high-frequency low-quality → low-frequency high-quality", 1),
    ], notes_text=(
        "Signal refinement is the primary lever for improving the gross return per fill. If the current "
        "signal fires indiscriminately in both high and low quality microstructure conditions, adding V2 "
        "features that select for high-persistence, high-confidence conditions should increase average "
        "quality of signal firings. The research question: do signals scoring highly on V2 features have "
        "meaningfully higher filled_avg_net? If yes, V2 thresholds convert the signal profile and "
        "directly improve NPA: fewer attempts, higher return per fill."
    ), visual_label=(
        "Side-by-side kernel density plots: left=distribution of filled_avg_net for all current signals. "
        "Right=distribution for signals surviving V2 filtering. "
        "Hypothesis: right distribution shifted rightward (higher per-fill returns). "
        "Label: 'TBD — pending V2 gate sweep'. Target: distributions separate with p-value < 0.05."
    ))

    _content_slide(prs, "9.2", "Gate Calibration — Finding the Optimal Threshold Vector", [
        ("Sweep: min_intensity_strong ∈ {2500, 3000, 3500, 4000, 5000}", 0),
        ("Sweep: min_imbalance_strong ∈ {0.5, 0.6, 0.7, 0.8}", 0),
        ("Sweep: max_spread_tight ∈ {0.0002, 0.0003, 0.0004, 0.0005}", 0),
        ("Total: 5 × 4 × 4 = 80 gate configurations per pocket — tractable (<30s per pocket)", 0),
        ("Objective: maximize net_per_attempt_after_gate on held-out validation splits", 0),
        ("Guard: gate_pass_rate > 0.15 — prevent over-restrictive gates (< 30 attempts = noisy estimate)", 0),
        ("Calibration discipline: sweep on training window, evaluate on validation window only", 1),
        ("Cross-validate gate configuration across multiple symbols for generalization evidence", 1),
    ], notes_text=(
        "Gate calibration is a hyperparameter optimization problem with clear objective (maximize NPA "
        "on held-out validation) and meaningful constraints (minimum attempt volume). 80 configurations "
        "is tractable because each evaluation is lightweight in-memory computation — feature data is "
        "cached. The minimum gate_pass_rate constraint (> 0.15) prevents optimizing NPA on a 10-attempt "
        "sample, which is noise optimization. Appropriate calibration protocol: training window for "
        "gate sweep, validation window for evaluation — same forward-split discipline as signal thresholds."
    ), visual_label=(
        "3D surface plot (or 2D heatmap with third dimension as annotation). "
        "x=min_intensity_strong, y=min_imbalance_strong, color=NPA at optimal max_spread_tight. "
        "Overlaid contour lines. Region where gate_pass_rate < 0.15 marked 'insufficient data'. "
        "Peak NPA point highlighted = target gate configuration."
    ))

    _content_slide(prs, "9.3", "Capacity Optimization and Ranking Sensitivity", [
        ("Capacity problem: high attempts_per_min at negative NPA = high-frequency loss generation", 0),
        ("Two remediation paths: (a) reduce attempt frequency via gate, (b) improve NPA directly", 0),
        ("Feature enrichment: order book depth snapshots, tick direction velocity, aggressor imbalance", 0),
        ("Alternative rules: search for rules with higher score_raw_core at reasonable attempt frequency", 0),
        ("Sensitivity analysis: bootstrap (seed, split) sets — measure ranking order stability", 0),
        ("Stable ranking: pocket maintains rank across 90%+ of bootstrap samples → trustworthy", 1),
        ("Unstable ranking: frequent swaps across samples → evaluation window too short, signal too noisy", 1),
    ], notes_text=(
        "Capacity and feature enrichment are medium-term directions requiring additional data "
        "infrastructure. The sensitivity analysis goal is methodologically important: if ranking order "
        "is highly sensitive to which 5 seeds or 4 splits are chosen, the ranking is not a reliable "
        "guide to live performance. Bootstrapping — generating many alternative (seed, split) sets "
        "and measuring ranking stability — quantifies this reliability. A stability score per pocket "
        "tells us how much to trust its rank."
    ), visual_label=(
        "Ranked list of top-10 pockets shown three times: current config, alternative seeds, "
        "alternative splits. Colored arrows show how rankings change (up=improvement, down=degradation). "
        "'Stability score' annotation per pocket: frequency of appearing in top-5 across configurations."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # SECTION 10 — PHILOSOPHICAL INSIGHT
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "X", "Philosophical and Systems Insight",
                     "The Deeper Principles That Govern the Design")

    _content_slide(prs, "10.1", "Markets Punish Naive Alpha — Execution Realism Determines Truth", [
        ("Apparent alpha = gross return on fills, measured without costs", 0),
        ("Real alpha = net return per attempt, after all execution costs, spread over all signals", 0),
        ("The gap between apparent and real is the cost of market participation", 0),
        ("In microstructure strategies, this gap frequently exceeds the apparent edge", 0),
        ("The purpose of execution realism is to measure this gap with maximum fidelity", 0),
        ("A 70% accurate signal requiring 3 bps of execution cost is no better than a coin flip at 2.5 bps gross", 1),
        ("Published signals that ignore fill selectivity and adverse selection are not deployable at scale", 1),
    ], notes_text=(
        "Markets do not reward clever signal identification; they reward the net of signal identification "
        "minus the cost of acting on it. The 'naive alpha' trap is well-documented in microstructure "
        "literature. Researchers identify predictable order flow patterns, compute gross returns, and "
        "conclude alpha exists. They fail to account for market impact, adverse selection, fees, and "
        "fill selectivity. Eclipse Scalper's architecture explicitly models the gap. Every component "
        "that appears to be operational overhead is doing epistemological work: measuring something "
        "real about the strategy's expected performance in a live environment."
    ), visual_label=(
        "Two-tier horizontal bar for a stylized strategy. "
        "Top: 'Backtesting without execution model' — large green bar labeled 'Apparent Alpha'. "
        "Bottom: same bar with cost components carved out: fee (red), adverse (amber), missed fills (grey). "
        "Remaining slice = 'Real Alpha'. For our system: carve-outs exceed apparent alpha → "
        "remaining slice is negative. This is the core finding visualized."
    ))

    _content_slide(prs, "10.2", "The System as Distributed Cognitive Architecture", [
        ("Signal (strategies layer) = Perception — what does the market state suggest?", 0),
        ("Gate (entry_loop + attempt gate) = Judgment — is now the right moment to act?", 0),
        ("Execution (order_router + exchange) = Action — commit to the decision faithfully", 0),
        ("Reconcile (reconcile.py) = Truth Correction — align belief with exchange reality", 0),
        ("Position manager = Consequence Management — handle the outcome mechanically", 0),
        ("Each layer irreplaceable; no shortcut from perception to action is safe", 0),
        ("Most value is added or destroyed at the Judgment layer — gate calibration is the key", 1),
        ("Execution layer must not second-guess Judgment — execute faithfully, report accurately", 1),
    ], notes_text=(
        "This framing reveals the system's design philosophy. The layers correspond to cognitive "
        "functions that any effective decision-making agent must perform. Perception is necessarily "
        "uncertain and approximate — the raw signal carries noise. Judgment is where most value is "
        "added or destroyed: premature action (weak signals) and delayed action (missing strong signals) "
        "both destroy value. The gate calibration problem is exactly the problem of making good judgment "
        "under uncertainty. Execution does not second-guess judgment — it executes faithfully. "
        "Reconciliation answers: 'What actually happened?' — the truth-correction mechanism."
    ), visual_label=(
        "Pentagon diagram with five vertices: Signal / Gate / Execution / Reconcile / Position_Manager. "
        "Directed edges connecting each to next. Color: blue=perception layers, green=action layers, "
        "red=truth layers. Each vertex annotated: function performed, key module, failure mode if absent."
    ))

    _content_slide(prs, "10.3", "From Signal Generator to Execution-Aware Organism", [
        ("v1 (initial): Signal → order. No gate. No capacity model. No execution realism.", 0),
        ("v2: Entry_loop gates. Kill-switch. Circuit breaker. Intent lifecycle. EXE invariants.", 0),
        ("v3: Passive fill simulator. Forward validation. Per-attempt metrics introduced.", 0),
        ("v4: Capacity filter. NPA-based ranking. Score raw fields. Gate degeneracy eliminated.", 0),
        ("v5 (current): Pre-attempt gate. DAT-04 verified. val_attempts_before_gate diagnostic.", 0),
        ("Each version reduces gap between apparent and real alpha measurement", 0),
        ("Trajectory is convergent: each version reduces illusion, increases fidelity", 1),
        ("v6 target: positive NPA confirmed on held-out splits at realistic fee/adverse", 1),
    ], notes_text=(
        "The evolutionary trajectory is a story of progressively better measurement fidelity. v1 was a "
        "signal generator dressed as a trading system. v2 added execution infrastructure. v3 first modeled "
        "execution reality — fill rates below 100% changed the research picture dramatically. v4 fixed "
        "the ranking system to reflect this. v5 adds the pre-attempt gate — first version with a "
        "research-validated path toward positive NPA. The trajectory is convergent: each version reduces "
        "illusion and increases accuracy. The goal is not to find alpha by any means — it is to find "
        "alpha that survives measurement."
    ), visual_label=(
        "Timeline with 5 version markers. Each: version number, key addition, measurement fidelity score "
        "(low → very high). Ascending curve overlaid showing fidelity improvement. "
        "Annotations at each version: gap closed. Dotted line extending beyond v5 "
        "toward 'v6: positive NPA confirmed'."
    ))

    _content_slide(prs, "10.4", "Operating Philosophy — Engineering Honesty in Quantitative Research", [
        ("Measure what you claim to measure — not a proxy that is easier to measure", 0),
        ("A negative result that is honest is more valuable than a positive result that is not", 0),
        ("Every claim must be backed by a reproducible test or empirical verification", 0),
        ("The pipeline's job is to generate falsifiable hypotheses, not to confirm priors", 0),
        ("Deployment readiness: NPA > 0 on held-out forward splits at fee=1.0, adv_mult=1.0", 0),
        ("'Measure more carefully, not deploy more aggressively' — the governing heuristic", 0),
        ("Precision of measurement is the competitive moat in microstructure research", 1),
        ("79 tests passing is a property of the current implementation — it must be maintained", 1),
    ], notes_text=(
        "This slide codifies the philosophical commitments that guided every architectural decision. "
        "'Measure what you claim to measure' — if we claim to measure expected PnL, we must measure NPA, "
        "not per-fill returns, not gross returns, not in-sample PnL. 'A negative result that is honest "
        "is more valuable than a positive result that is not' — the system currently shows negative NPA. "
        "This is correct and valuable. An alternative system showing positive NPA by using in-sample "
        "data or ignoring fill rates would be showing a dishonest positive — deploying on it causes losses. "
        "Deployment trigger is concrete and falsifiable: NPA > 0 at fee=1.0, adv=1.0 on held-out splits."
    ), visual_label=(
        "Minimalist three-column table: Commitment | Mechanism | Verification. "
        "Row 1: Measure correctly | NPA not per-fill | test_npa_cost_applied_to_fills_only. "
        "Row 2: Honest negatives | NPA gate zeroes non-viable | test_ranker_filters_low_capacity. "
        "Row 3: Reproducible | seed-based determinism | test_net_per_attempt_deterministic. "
        "Row 4: Deployment gate | NPA>0 criterion | Forward validation pipeline. "
        "Footer: '79 tests passing as of 2026-02-20'."
    ))

    # ══════════════════════════════════════════════════════════════════════
    # CLOSING
    # ══════════════════════════════════════════════════════════════════════
    _closing_slide(prs)

    # ══════════════════════════════════════════════════════════════════════
    # APPENDICES
    # ══════════════════════════════════════════════════════════════════════
    _section_divider(prs, "⊕", "Appendices",
                     "File Map · Observability · Repro Commands")

    _summary_slide(prs, "A.1", "Appendix A — File-to-Concern Map", [
        ("execution/entry_loop.py",             "Execution Gate",  "20+ gate checks before submission; kill-switch, circuit breaker"),
        ("execution/order_router.py",           "Execution Route", "Idempotent routing, retry classification, clientOrderId"),
        ("execution/reconcile.py",              "State Authority", "Exchange→brain state convergence; orphan detection"),
        ("execution/passive_execution_simulator.py", "Research Sim", "simulate_passive_fill: deterministic, calibrated cost model"),
        ("tools/micro_edge_backtest.py",        "Research Backtest","simulate_rule_trades + _passes_attempt_gate + metrics"),
        ("tools/validate_passive_pocket_forward.py", "Research Val","Forward splits, capacity fields, gate diagnostic outputs"),
        ("tools/rank_passive_pockets_forward.py","Research Rank",  "NPA-based scoring, capacity filter, score_raw_* fields"),
        ("tools/micro_edge_lib.py",             "Feature Gen",     "build_bucket_features: 1-second bucket aggregation"),
        ("tests/test_exec_cost_models.py",      "Test: Costs",     "test_npa_cost_applied_to_fills_only: DAT-04 verification"),
        ("tests/test_passive_realistic_sim.py", "Test: Gate",      "test_attempt_gate_reduces_n_attempts: gate structural test"),
    ], notes_text="Navigation aid for engineers exploring the codebase.")

    _content_slide(prs, "A.2", "Appendix B — Observability Conventions", [
        ("All telemetry: JSONL format, append-only in logs/telemetry.jsonl", 0),
        ("Log prefix: [DEBUG_FLOW source=module:function] — grep-able by module path", 0),
        ("Research reports: markdown in reports/, machine-readable JSON in reports/*.json", 0),
        ("Invariant violations: [INVARIANT_VIOLATION invariant=EXE-01] prefix — alert on this pattern", 0),
        ("Gate diagnostic: debug_stats['attempt_gate_blocked'] — accessible in every sim result dict", 0),
        ("Per-combo: val_attempts_before_gate vs val_attempts_after_gate — gate impact visible per row", 0),
        ("Test naming: test_{what_is_verified}_{condition} — descriptive and searchable", 1),
    ], notes_text=(
        "Observability conventions make the system's behavior inspectable without modifying it. "
        "JSONL telemetry allows offline analysis with jq, pandas, custom dashboards. Structured log "
        "prefixes enable grep-based filtering. Per-combo gate diagnostic fields allow direct inspection "
        "of gate effect at the individual (seed, split) level without aggregation loss."
    ), has_visual=False)

    _code_slide(prs, "A.3", "Appendix C — Research Repro Commands", [
        "# ── Forward validate single pocket with gate ──────────────────────",
        "python -m tools.validate_passive_pocket_forward \\",
        "  --db data/microstructure.db --symbol ETHUSDT \\",
        "  --lookback-min 1440 --bucket-sec 1 --horizon-sec 120 \\",
        "  --rule intensity_spike_imbalance_cont --side auto \\",
        "  --min-imbalance 0.5 --min-trade-intensity 2500 --max-spread 0.0005 \\",
        "  --maker-fee-bps 1.0 \\",
        "  --min-intensity-strong 3500 --max-spread-tight 0.0003 \\",
        "  --splits 4 --seeds 11,22,33,44,55 --min-n 30",
        "",
        "# ── Rank pockets with capacity filter and gate ────────────────────",
        "python -m tools.rank_passive_pockets_forward \\",
        "  --candidates-md reports/FILTER_SWEEP_PASSIVE_REALISTIC_ETH.md \\",
        "  --db data/microstructure.db \\",
        "  --maker-fee-bps-grid 0.5,1.0,1.5 \\",
        "  --passive-adverse-mult-grid 0.8,1.0,1.2 \\",
        "  --min-attempt-fill-rate 0.10 \\",
        "  --min-intensity-strong 3500 --max-spread-tight 0.0003 \\",
        "  --out-md reports/PASSIVE_POCKET_RANKING.md \\",
        "  --out-json reports/PASSIVE_POCKET_RANKING.json",
        "",
        "# ── Run full test suite ───────────────────────────────────────────",
        "python -m pytest -q   # → 79 passed",
    ], notes_text=(
        "These commands are the canonical repro steps for the current research state. "
        "Parameterized with values used to generate the results in this presentation. "
        "Running against the same database produces identical results (DAT-03)."
    ))


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def _args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Generate Eclipse Scalper technical PowerPoint deck.")
    p.add_argument("--out", default="reports/eclipse_scalper_technical_deck.pptx",
                   help="Output .pptx path")
    return p.parse_args()


def main() -> int:
    args = _args()
    out = Path(str(args.out))
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    build_all_slides(prs)

    n = len(prs.slides)
    print(f"  {n} slides generated.")

    prs.save(str(out))
    print(f"Saved -> {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
